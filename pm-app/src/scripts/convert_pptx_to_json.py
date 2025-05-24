import os
import json
import zipfile
import xml.etree.ElementTree as ET
import uuid
from pptx import Presentation
from pptx.text.fonts import FontFiles

def get_all_fonts_from_theme(pptx_path):
    """Lấy tất cả font từ theme1.xml, bao gồm major_font, minor_font và các font theo script."""
    try:
        with zipfile.ZipFile(pptx_path, 'r') as zip_ref:
            with zip_ref.open('ppt/theme/theme1.xml') as theme_file:
                tree = ET.parse(theme_file)
                root = tree.getroot()
                ns = {'a': 'http://schemas.openxmlformats.org/drawingml/2006/main'}
                
                # Lấy major_font
                major_fonts = {}
                for major_font_elem in root.findall('.//a:majorFont/a:latin', ns):
                    major_font = major_font_elem.get('typeface') 
                    if major_font:
                        major_fonts = major_font
                    else: None

                
                # Lấy minor_font
                minor_fonts = {}
                for minor_font_elem in root.findall('.//a:minorFont/a:latin', ns):
                    minor_font = minor_font_elem.get('typeface')
                    if minor_font:
                        minor_fonts = major_font
                    else: None
                
                print(f"  Debug - Major font: {major_font}")
                print(f"  Debug - Minor font: {minor_font}")
                
                return {
                    "major_font": major_font,
                    "minor_font": minor_font
                }
    except Exception as e:
        print(f"  Debug - Error reading theme from {pptx_path}: {e}")
        return None
    
def get_font_and_size_from_slide_xml(pptx_path, slide_num=1):
    """Lấy font, size và line height từ slideX.xml."""
    try:
        with zipfile.ZipFile(pptx_path, 'r') as zip_ref:
            slide_path = f'ppt/slides/slide{slide_num}.xml'
            with zip_ref.open(slide_path) as slide_file:
                tree = ET.parse(slide_file)
                root = tree.getroot()
                ns = {'a': 'http://schemas.openxmlformats.org/drawingml/2006/main', 'p': 'http://schemas.openxmlformats.org/presentationml/2006/main'}
                
                fonts_and_sizes = {
                    "title": {"font": None, "size": None, "line_height": None},
                    "content": {"font": None, "size": None, "line_height": None}
                }
                
                # Tìm các shape dựa trên placeholder type
                for sp in root.findall('.//p:sp', ns):
                    ph = sp.find('.//p:ph', ns)
                    if ph is not None:
                        is_title = ph.get('type') == 'title'
                        is_content = ph.get('type') == 'body' or ph.get('idx') == '1'
                        
                        # Lấy font, size và line height từ <a:rPr> và <a:pPr>
                        run_props = sp.findall('.//a:rPr', ns)
                        for rpr in run_props:
                            # Font từ <a:latin> hoặc <a:buFont>
                            latin = rpr.find('a:latin', ns)
                            if latin is not None and latin.get('typeface'):
                                font = latin.get('typeface')
                            else:
                                bu_font = rpr.find('a:buFont', ns)
                                font = bu_font.get('typeface') if bu_font and bu_font.get('typeface') else None
                            
                            # Size từ sz (đơn vị: 1/100 pt)
                            size = rpr.get('sz')
                            size = int(size) / 100 if size else None
                            
                            # Line height từ <a:pPr><a:lnSpc>
                            p_pr = sp.find('.//a:pPr', ns)
                            if p_pr is not None:
                                ln_spc = p_pr.find('a:lnSpc', ns)
                                if ln_spc is not None:
                                    spc_pct = ln_spc.find('a:spcPct', ns)
                                    spc_pts = ln_spc.find('a:spcPts', ns)
                                    if spc_pct is not None:
                                        line_height = int(spc_pct.get('val')) / 100000  # Chuyển đổi sang phần trăm
                                    elif spc_pts is not None:
                                        line_height = int(spc_pts.get('val')) / 100  # Chuyển đổi sang điểm
                                    else:
                                        line_height = 100  # Mặc định single line spacing
                                else:
                                    line_height = 100  # Mặc định nếu không có lnSpc
                            else:
                                line_height = 100  # Mặc định nếu không có pPr
                            
                            if font or size or line_height:
                                if is_title:
                                    fonts_and_sizes["title"]["font"] = font
                                    fonts_and_sizes["title"]["size"] = size
                                    fonts_and_sizes["title"]["line_height"] = line_height
                                elif is_content:
                                    fonts_and_sizes["content"]["font"] = font
                                    fonts_and_sizes["content"]["size"] = size
                                    fonts_and_sizes["content"]["line_height"] = line_height
                        
                        # Kiểm tra <a:endParaRPr> nếu không tìm thấy size
                        end_para_rpr = sp.find('.//a:endParaRPr', ns)
                        if end_para_rpr is not None:
                            size = end_para_rpr.get('sz')
                            size = int(size) / 100 if size else None
                            if is_title and not fonts_and_sizes["title"]["size"]:
                                fonts_and_sizes["title"]["size"] = size
                            elif is_content and not fonts_and_sizes["content"]["size"]:
                                fonts_and_sizes["content"]["size"] = size
                
                print(f"  Debug - Fonts, sizes, and line heights from slide{slide_num}.xml: {fonts_and_sizes}")
                return fonts_and_sizes
    except Exception as e:
        print(f"  Debug - Error reading slide{slide_num}.xml: {e}")
        return None

def get_font_info(paragraph, shape=None, slide=None, prs=None, pptx_path=None):
    # 1. Lấy từ run trong shape (python-pptx)
    for run in paragraph.runs:
        if run.text.strip():
            font_name = run.font.name
            font_size = int(run.font.size.pt) if run.font.size else None
            # python-pptx không hỗ trợ lấy line height trực tiếp, bỏ qua ở đây
            print(f"  Debug - Run text: '{run.text}', Font: {font_name}, Size: {font_size}")
            if font_name or font_size:
                return font_name, font_size, None  # Trả về None cho line_height
    
    # 2. Lấy từ paragraph trong shape (python-pptx)
    font_name = paragraph.font.name if paragraph.font and paragraph.font.name else None
    font_size = int(paragraph.font.size.pt) if paragraph.font and paragraph.font.size else None
    print(f"  Debug - Paragraph text: '{paragraph.text.strip()}', Font: {font_name}, Size: {font_size}")
    if font_name:
        return font_name, font_size, None  # Trả về None cho line_height
    
    # 3. Lấy từ slide XML nếu python-pptx không tìm thấy
    if pptx_path and slide:
        slide_num = prs.slides.index(slide) + 1
        slide_data = get_font_and_size_from_slide_xml(pptx_path, slide_num)
        if slide_data:
            is_title = shape and shape.is_placeholder and shape.placeholder_format.type == 1  # PP_PLACEHOLDER.TITLE
            font_name = slide_data["title"]["font"] if is_title else slide_data["content"]["font"]
            font_size = slide_data["title"]["size"] if is_title else slide_data["content"]["size"]
            line_height = slide_data["title"]["line_height"] if is_title else slide_data["content"]["line_height"]
            if font_name or font_size or line_height:
                print(f"  Debug - Fallback to slide{slide_num}.xml: Font: {font_name}, Size: {font_size}, Line Height: {line_height}")
                return font_name, font_size, line_height
    
    # 4. Lấy từ theme nếu vẫn không tìm thấy
    if pptx_path:
        theme_fonts = get_all_fonts_from_theme(pptx_path)
        if theme_fonts:
            is_title = shape and shape.is_placeholder and shape.placeholder_format.type == 1  # PP_PLACEHOLDER.TITLE
            font_name = theme_fonts["major_font"] if is_title else theme_fonts["minor_font"]
            font_size = 30 if is_title else 115  # Giá trị mặc định nếu không có size
            line_height = 100  # Mặc định single line spacing
            print(f"  Debug - Fallback to theme: Font: {font_name}, Size: {font_size}, Line Height: {line_height}")
            return font_name, font_size, line_height
    
    print(f"  Debug - No font, size, or line height found for '{paragraph.text.strip()}'")
    return None, None, None

def pptx_folder_to_json(folder_path, output_file):
    fonts = FontFiles._installed_fonts()
    font_names = [i[0] for i in fonts]
    print("Installed fonts on system:", font_names)
    
    if not os.path.exists(folder_path):
        print(f"Folder {folder_path} does not exist")
        return
    
    output_dir = os.path.dirname(output_file)
    if not os.path.exists(output_dir):
        os.makedirs(output_dir)
    
    all_files = []
    file_id = 0
    
    for file_name in os.listdir(folder_path):
        if file_name.endswith(".pptx"):
            file_path = os.path.join(folder_path, file_name)
            try:
                prs = Presentation(file_path)
            except Exception as e:
                print(f"Error opening {file_name}: {e}")
                continue
            
            theme_fonts = get_all_fonts_from_theme(file_path)
            if not theme_fonts:
                continue
            
            slides = []
            for i, slide in enumerate(prs.slides):
                print(f"\nProcessing Slide {i+1}")
                title = None
                content = None
                title_font = None
                title_font_size = None
                content_font = None
                content_font_size = None
                content_line_height = None
                
                for shape in slide.shapes:
                    if not hasattr(shape, "text_frame") or not shape.text_frame:
                        continue
                    text = shape.text_frame.text.strip()
                    if not text:
                        continue
                    
                    print(f"Shape text: '{text}'")
                    for paragraph in shape.text_frame.paragraphs:
                        paragraph_text = paragraph.text.strip()
                        if not paragraph_text:
                            continue
                        
                        font_name, font_size, line_height = get_font_info(paragraph, shape, slide, prs, pptx_path=file_path)
                        
                        is_title_shape = shape.is_placeholder and shape.placeholder_format.type == 1  # PP_PLACEHOLDER.TITLE
                        is_short_text = len(paragraph_text) < 50
                        
                        if is_title_shape or (not title and is_short_text and font_size and font_size <= 40):
                            if not title:
                                title = paragraph_text
                                title_font = font_name
                                title_font_size = font_size
                        else:
                            if not content or content == title:
                                content = paragraph_text
                                content_font = font_name
                                content_font_size = font_size
                                content_line_height = line_height
                
                # Gán font, size và line height từ theme nếu vẫn null
                if content:
                    if not title_font and theme_fonts["major_font"]:
                        title_font = theme_fonts["major_font"]
                        title_font_size = title_font_size or 30
                    if not content_font and theme_fonts["minor_font"]:
                        content_font = theme_fonts["minor_font"]
                        content_font_size = content_font_size or 115
                        content_line_height = content_line_height or 100
                    
                    slide_data = {
                        "id": str(uuid.uuid4()),
                        "title": {
                            "text": title or "",
                            "fontName": title_font or "Unknown",
                            "fontSize": title_font_size or 30
                        },
                        "content": {
                            "text": content or "",
                            "fontName": content_font or "Unknown",
                            "fontSize": content_font_size or 115,
                            # "lineHeight": content_line_height or 100  # Thêm lineHeight
                        }
                    }
                    slides.append(slide_data)
                    print(f"Added slide: {slide_data}")

            # Loại bỏ đuôi .pptx
            base_file_name = os.path.splitext(file_name)[0]
            
            file_data = {
                "id": str(uuid.uuid4()),
                "fileName": base_file_name,
                "slides": slides,
            }
            all_files.append(file_data)
            file_id += 1
    
    if all_files:
        with open(output_file, "w", encoding="utf-8") as json_file:
            json.dump(all_files, json_file, ensure_ascii=False, indent=4)
        print(f"All data saved to: {output_file}")
    else:
        print(f"No data extracted from {folder_path}")

# Chạy chương trình
folder_path = r'../data/pptx'
output_file = r"../data/output/output.json"
pptx_folder_to_json(folder_path, output_file)