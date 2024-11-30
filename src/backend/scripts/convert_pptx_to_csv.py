import os
import pandas as pd
from pptx import Presentation

def extract_song_data(pptx_path):
    """
    Trích xuất dữ liệu bài hát từ một tệp PPTX.

    Tham số:
        pptx_path (str): Đường dẫn đến tệp PPTX.

    Trả về:
        dict: Một dictionary với các khóa 'Song_name' và 'Song_lyrics'.
    """
    try:
        prs = Presentation(pptx_path)
    except Exception as e:
        print(f"Error opening {pptx_path}: {e}")
        return None

    song_data = {
        'Song_name': '',
        'Song_lyrics': ''
    }

    lyrics_parts = []

    for slide in prs.slides:
        # Trích xuất tiêu đề (Song_name)
        title = slide.shapes.title
        if title and title.text:
            song_data['Song_name'] = title.text.strip()

        # Trích xuất nội dung (Song_lyrics)
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            text = shape.text_frame.text.strip()
            # Bỏ qua văn bản tiêu đề
            if text == song_data['Song_name']:
                continue
            # Thêm văn bản vào lyrics_parts
            lyrics_parts.append(text)

    # Kết hợp tất cả phần của lời bài hát
    song_data['Song_lyrics'] = '\n'.join(lyrics_parts).strip()

    return song_data

def main():
    # Định nghĩa đường dẫn
    project_root = os.getcwd()
    pptx_dir = os.path.join(project_root, 'src', 'backend', 'pptx-generator', 'data', 'pptx')
    output_dir = os.path.join(project_root, 'src', 'backend', 'pptx-generator', 'data', 'output')
    output_csv = os.path.join(output_dir, 'songs.csv')

    print(f"Project Root: {project_root}")
    print(f"PPTX Directory: {pptx_dir}")
    print(f"Output Directory: {output_dir}")
    print(f"Output CSV: {output_csv}")

    # Đảm bảo thư mục output tồn tại
    os.makedirs(output_dir, exist_ok=True)

    # Khởi tạo danh sách lưu trữ dữ liệu bài hát
    all_songs = []

    # Kiểm tra xem thư mục pptx_dir có tồn tại không
    if not os.path.exists(pptx_dir):
        print(f"Error: The directory {pptx_dir} does not exist.")
        return

    # Liệt kê tất cả các tệp PPTX
    pptx_files = [f for f in os.listdir(pptx_dir) if f.lower().endswith('.pptx')]
    if not pptx_files:
        print(f"No PPTX files found in {pptx_dir}.")
        return

    print(f"Found {len(pptx_files)} PPTX files to process.")

    # Duyệt qua tất cả các tệp PPTX trong pptx_dir
    for filename in pptx_files:
        pptx_path = os.path.join(pptx_dir, filename)
        print(f"Processing file: {filename}")
        song = extract_song_data(pptx_path)
        if song:
            if song['Song_name']:
                all_songs.append(song)
                print(f"Extracted: {song['Song_name']}")
            else:
                print(f"Warning: 'Song_name' not found in {filename}")
        else:
            print(f"Failed to extract data from {filename}")

    # Kiểm tra xem có bài hát nào được trích xuất không
    if not all_songs:
        print("No song data was extracted. Please check the PPTX files' structure.")
        return

    # Tạo DataFrame
    df = pd.DataFrame(all_songs, columns=['Song_name', 'Song_lyrics'])

    # Lưu vào CSV
    try:
        df.to_csv(output_csv, index=False, encoding='utf-8-sig')  # 'utf-8-sig' để hiển thị đúng trong Excel
        print(f"\nAll data has been successfully saved to {output_csv}")
    except Exception as e:
        print(f"Error saving CSV file: {e}")

if __name__ == "__main__":
    main()
